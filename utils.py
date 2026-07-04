from __future__ import annotations

from io import BytesIO
from typing import Any

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def generate_word_document(data: dict[str, Any]) -> bytes:
    """Generate an executive Word briefing from validated presales data."""
    document = Document()
    document.add_heading("AI Delivery & Presales Executive Briefing", 0)
    document.add_paragraph(f"Generated: {data.get('generated_at', 'N/A')}")

    scope = data["scope_requirements"]
    document.add_heading("Scope & Functional Requirements", level=1)
    document.add_paragraph(scope.get("executive_summary", ""))

    add_bulleted_section(document, "Objectives", scope.get("objectives", []))
    add_bulleted_section(document, "Deliverables", scope.get("deliverables", []))
    add_bulleted_section(document, "Core Constraints", scope.get("core_constraints", []))

    document.add_heading("Chronological Deadlines", level=2)
    for item in scope.get("chronological_deadlines", []):
        document.add_paragraph(
            f"{item.get('timeframe', 'TBD')}: {item.get('milestone', 'Milestone')} - "
            f"{item.get('business_value', '')}",
            style="List Bullet",
        )

    document.add_heading("RAID Log Matrix", level=1)
    risks = data.get("raid_risk_matrix", [])
    table = document.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    headers = ["Risk Factor", "Category", "Impact Severity", "Proactive Mitigation Strategy"]
    for index, header in enumerate(headers):
        table.rows[0].cells[index].text = header
    for risk in risks:
        row = table.add_row().cells
        for index, header in enumerate(headers):
            row[index].text = str(risk.get(header, ""))

    qa = data["qa_intelligence"]
    document.add_heading("QA Intelligence & Staffing", level=1)
    document.add_paragraph(f"Recommended Validation Approach: {qa.get('recommended_testing_stack', '')}")
    document.add_paragraph(f"Staffing Allocation Index: {qa.get('staffing_allocation_index', 'N/A')} FTE")
    add_bulleted_section(document, "Testing Strategy", qa.get("testing_strategy", []))
    add_bulleted_section(document, "Tooling Frameworks", qa.get("tooling_frameworks", []))

    document.add_heading("Staffing Matrix", level=2)
    for person in qa.get("staffing_matrix", []):
        document.add_paragraph(
            f"{person.get('allocation', '')} {person.get('role', '')}: {person.get('responsibility', '')}",
            style="List Bullet",
        )

    output = BytesIO()
    document.save(output)
    return output.getvalue()


def generate_presentation_slides(data: dict[str, Any]) -> bytes:
    """Generate a compact enterprise pitch deck from validated presales data."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    add_title_slide(
        presentation,
        "AI Delivery & Presales Briefing",
        "Executive delivery intelligence",
    )

    scope = data["scope_requirements"]
    add_bullet_slide(
        presentation,
        "Scope & Functional Requirements",
        [scope.get("executive_summary", "")] + scope.get("objectives", [])[:4],
    )
    add_bullet_slide(
        presentation,
        "Delivery Constraints & Deadlines",
        scope.get("core_constraints", [])[:4]
        + [
            f"{item.get('timeframe', 'TBD')}: {item.get('milestone', '')}"
            for item in scope.get("chronological_deadlines", [])[:4]
        ],
    )
    add_bullet_slide(
        presentation,
        "RAID Matrix Highlights",
        [
            f"{risk.get('Impact Severity', 'Medium')} | {risk.get('Category', '')}: {risk.get('Risk Factor', '')}"
            for risk in data.get("raid_risk_matrix", [])[:5]
        ],
    )

    qa = data["qa_intelligence"]
    add_bullet_slide(
        presentation,
        "QA Intelligence & Staffing",
        [
            f"Validation Approach: {qa.get('recommended_testing_stack', '')}",
            f"Staffing Index: {qa.get('staffing_allocation_index', 'N/A')} FTE",
        ]
        + [
            f"{person.get('allocation', '')} {person.get('role', '')}"
            for person in qa.get("staffing_matrix", [])[:4]
        ],
    )

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def generate_pdf_document(data: dict[str, Any]) -> bytes:
    """Reserved PDF export framework for production expansion."""
    output = BytesIO()
    pdf = canvas.Canvas(output, pagesize=letter)
    pdf.setTitle("AI Delivery & Presales Executive Briefing")
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(72, 740, "AI Delivery & Presales Executive Briefing")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(72, 720, f"Generated: {data.get('generated_at', 'N/A')}")
    pdf.drawString(72, 704, "PDF export framework is wired for production extension.")
    pdf.save()
    return output.getvalue()


def add_bulleted_section(document: Any, title: str, items: list[str]) -> None:
    document.add_heading(title, level=2)
    for item in items:
        document.add_paragraph(str(item), style="List Bullet")


def add_title_slide(presentation: Any, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(presentation: Any, title: str, items: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.level = 0
        paragraph.font.size = Pt(18)
